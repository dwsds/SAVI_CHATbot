
import os
from docx import Document
from pptx import Presentation
import fitz 

def extract_text_from_docx(path):
    """Extract text from DOCX files"""
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_pptx(path):
    """Extract text from PPTX files"""
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(path):
    """Extract text from PDF files"""
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def load_all_data(base_folder):
    """Load and extract text from supported file types"""
    corpus = ""
    
    if not os.path.exists(base_folder):
        os.makedirs(base_folder)
        print(f"Created {base_folder} directory. Please add your documents there.")
        return corpus
    
    file_count = 0
    for root, _, files in os.walk(base_folder):
        for file in files:
            path = os.path.join(root, file)
            try:
                if file.endswith(".docx"):
                    corpus += extract_text_from_docx(path) + "\n"
                    file_count += 1
                elif file.endswith(".pptx"):
                    corpus += extract_text_from_pptx(path) + "\n"
                    file_count += 1
                elif file.endswith(".pdf"):
                    corpus += extract_text_from_pdf(path) + "\n"
                    file_count += 1
                elif file.endswith(".txt"):
                    with open(path, encoding="utf8") as f:
                        corpus += f.read() + "\n"
                    file_count += 1
            except Exception as e:
                print(f"Error processing {file}: {str(e)}")
    
    print(f"Processed {file_count} files")
    return corpus

text_data = load_all_data("data/")
print(f"{len(text_data)} characters loaded from documents.")

if len(text_data) == 0:
    print("WARNING: No text data loaded! Add documents to data/ folder.")

